#!/usr/bin/env python3
# 발표용 PPT 생성 (심플 + 표 위주)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x2B, 0x36, 0x5E); INK = RGBColor(0x22, 0x26, 0x2E)
GRAY = RGBColor(0x5C, 0x61, 0x6B); ACC = RGBColor(0xD2, 0x3A, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF); ROW = RGBColor(0xF2, 0xF4, 0xF7)
SUB = RGBColor(0x9F, 0xAC, 0xCF)
GOOD = RGBColor(0x1F, 0x7A, 0x4D)
STEPBG = RGBColor(0xF2, 0xF4, 0xF7)
FONT = "Apple SD Gothic Neo"
W, H = Inches(13.333), Inches(7.5)

prs = Presentation(); prs.slide_width = W; prs.slide_height = H
blank = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(blank)


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h); tb.text_frame.word_wrap = True
    return tb.text_frame


def setp(p, txt, size, color=INK, bold=False):
    p.text = txt
    for r in p.runs:
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = FONT


def addp(tf, txt, size, color=INK, bold=False, bullet=False, space=6):
    p = tf.add_paragraph()
    setp(p, ("•  " if bullet else "") + txt, size, color, bold)
    p.space_after = Pt(space)
    return p


def title_bar(s, kicker, title):
    bar = s.shapes.add_shape(1, 0, 0, W, Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    tf = box(s, Inches(0.6), Inches(0.45), Inches(12.1), Inches(1.1))
    setp(tf.paragraphs[0], kicker, 13, ACC, True)
    addp(tf, title, 30, NAVY, True, space=0)


def table(s, rows, l, t, w, colw=None, fs=14):
    nr, nc = len(rows), len(rows[0])
    gt = s.shapes.add_table(nr, nc, l, t, w, Inches(0.42 * nr)).table
    if colw:
        for i, cw in enumerate(colw):
            gt.columns[i].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = gt.cell(ri, ci); c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = Inches(0.12); c.margin_right = Inches(0.08)
            c.margin_top = Inches(0.03); c.margin_bottom = Inches(0.03)
            p = c.text_frame.paragraphs[0]; setp(p, str(val), fs)
            if ri == 0:
                c.fill.solid(); c.fill.fore_color.rgb = NAVY
                for r in p.runs:
                    r.font.color.rgb = WHITE; r.font.bold = True
            else:
                c.fill.solid(); c.fill.fore_color.rgb = ROW if ri % 2 else WHITE
    return gt


def step_box(s, x, y, w, h, num, header, lines, hcolor):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = STEPBG
    shp.line.color.rgb = hcolor; shp.line.width = Pt(2)
    shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.22); tf.margin_bottom = Inches(0.15)
    setp(tf.paragraphs[0], f"{num}  {header}", 18, hcolor, True)
    tf.paragraphs[0].space_after = Pt(10)
    for ln in lines:
        addp(tf, ln, 13, INK, bullet=True, space=8)
    return shp


def arrow(s, x, y):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, Inches(0.6), Inches(0.5))
    a.fill.solid(); a.fill.fore_color.rgb = NAVY; a.line.fill.background()
    a.shadow.inherit = False


# 1. 표지
s = slide()
bg = s.shapes.add_shape(1, 0, 0, W, H); bg.fill.solid()
bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
tf = box(s, Inches(1), Inches(2.4), Inches(11.3), Inches(2.8))
setp(tf.paragraphs[0], "인스타그램 콘텐츠 분석 자동화", 40, WHITE, True)
addp(tf, "게시물 성과 지표 수집 + AI 기반 주제 요약", 22, RGBColor(0xC9, 0xD2, 0xEA), space=18)
addp(tf, "사진 속 번체중문을 읽어 한국어 제목을 자동 생성하는 VLM 선정기", 15, SUB)
addp(tf, "VEASLY · 대만 타겟 한국 상품 대리구매", 13, SUB, space=2)

# 2. 배경 — 시작 → 문제 → 해결 (도형 3개)
s = slide(); title_bar(s, "WHY", "왜 이걸 만들었나")
bw, bh, by = Inches(3.7), Inches(3.6), Inches(2.4)
x1, x2, x3 = Inches(0.7), Inches(4.8), Inches(8.9)
step_box(s, x1, by, bw, bh, "1", "시작",
         ["인스타 게시물 성과 지표(도달·조회·좋아요·저장)를 자동 수집·정리하는 게 원래 목표"], NAVY)
step_box(s, x2, by, bw, bh, "2", "문제",
         ["지표는 모이는데 '무슨 내용인지' 한눈에 안 보임",
          "인스타 API는 제목/주제를 안 줌 (캡션 원문만)",
          "핵심은 사진 속 번체중문인데 데이터로 안 옴"], ACC)
step_box(s, x3, by, bw, bh, "3", "해결 방법",
         ["사진 + 캡션을 AI(VLM)로 분석",
          "→ 제목·주제 자동 생성",
          "지표 옆에 '무슨 글인지' 한눈에"], GOOD)
arrow(s, Inches(4.15), Inches(3.95))
arrow(s, Inches(8.25), Inches(3.95))

# 3. 1단계 후보 6개
s = slide(); title_bar(s, "STEP 1", "후보 모델 선정 — VLM 6개")
tf = box(s, Inches(0.7), Inches(1.7), Inches(12), Inches(0.8))
setp(tf.paragraphs[0], "선정 기준: 사진+텍스트 모델(VLM) · 로컬 24GB 실행(mlx) · 번체/한국어 · 인기 검증", 13, GRAY)
table(s, [["모델", "회사", "크기"],
          ["Qwen3-VL-8B", "알리바바", "8B"], ["Qwen3.5-9B", "알리바바", "9B"],
          ["InternVL3-8B", "상하이 AI Lab", "8B"], ["Gemma-4-12B", "구글", "12B"],
          ["Gemma-4-26B-A4B", "구글", "26B (MoE)"], ["Qwen3.6-27B", "알리바바", "27B"]],
      Inches(0.7), Inches(2.45), Inches(7.5), [Inches(3.2), Inches(2.6), Inches(1.7)])
tf = box(s, Inches(8.6), Inches(2.6), Inches(4), Inches(3))
setp(tf.paragraphs[0], "→ 서로 다른 회사 3곳", 15, NAVY, True)
addp(tf, "알리바바 · 구글 · 상하이AI", 13, GRAY, space=12)
addp(tf, "→ 크기 8~27B 다양", 15, NAVY, True)
addp(tf, "비교 다양성 확보", 13, GRAY)

# 4. 2단계
s = slide(); title_bar(s, "STEP 2", "6개 1차 비교 → 두 가지 발견")
tf = box(s, Inches(0.7), Inches(1.65), Inches(12), Inches(1.3))
setp(tf.paragraphs[0], "발견 ①  결과가 거의 다 'VEASLY 구매대행…'으로 비슷 → 게시물 구분 불가", 16, ACC, True)
addp(tf, "원인: 모든 게시물 공통 홍보문구를 그대로 가져옴  |  근거: 제목 60개 중 55개(92%)에 'VEASLY/구매' 포함", 13, GRAY)
table(s, [["모델", "결과 (같은 게시물)"],
          ["Qwen3.5-9B", "VEASLY 소맥피 남자 아이돌 대판점"],
          ["Gemma-26B", "VEASLY 한국 구매대행 아이돌 굿즈…"],
          ["InternVL3", "VEASLY 한국 대구매"],
          ["Qwen3-VL-8B", "VEASLY 한국 대행 소울푸드 남자 아이돌…"]],
      Inches(0.7), Inches(3.05), Inches(8.6), [Inches(2.6), Inches(6.0)], fs=13)
tf = box(s, Inches(0.7), Inches(5.85), Inches(12), Inches(1))
setp(tf.paragraphs[0], "발견 ②  Qwen3.6-27B → 24GB 초과로 실행 중 시스템 다운 → 사양 탈락 (남은 5개)", 16, ACC, True)

# 5. 3단계 프롬프트
s = slide(); title_bar(s, "STEP 3", "프롬프트 정교화")
tf = box(s, Inches(0.7), Inches(1.9), Inches(12), Inches(5))
setp(tf.paragraphs[0], "1.  입력 간소화 (시간 단축)", 17, NAVY, True)
addp(tf, "게시물의 여러 사진 → '첫 표지 사진 1장'만 사용", 14, INK, bullet=True, space=16)
addp(tf, "2.  캡션의 공통 홍보문구 배제 (결과 동일화 방지)", 17, NAVY, True)
addp(tf, "캡션은 입력하되, 모든 게시물 공통 문구는 제목에 쓰지 않도록 지시", 14, INK, bullet=True, space=16)
addp(tf, "3.  해시태그 활용 추가", 17, NAVY, True)
addp(tf, "고유 해시태그(#cortis 등)를 핵심 단서로, 일반 태그(#韓國代購)는 무시", 14, INK, bullet=True, space=16)
addp(tf, "→ 게시물별 고유 제목 포착   예) \"CORTIS W매거진 화보\"", 15, ACC, True)

# 6. 순서 실험 + 2개 탈락
s = slide(); title_bar(s, "STEP 3", "처리 순서 실험 (A/B) + 2개 탈락")
tf = box(s, Inches(0.7), Inches(1.65), Inches(12), Inches(0.8))
setp(tf.paragraphs[0], "OCR(번체)까지 동일 → A: 요약→번역  vs  B: 번역→요약    A 채택: 고유명사 보존", 14, GRAY)
table(s, [["게시물", "A (요약→번역) ✅", "B (번역→요약) ❌"],
          ["진기주(인명)", "진기주 전직 경력: 도전하라", "퇴사자들의 새로운 도전 (이름 증발)"],
          ["다이아몬드", "다이아몬드 신제품 (맥락 유지)", "입술용 마스카라? (환각)"],
          ["이준영(인명)", "이준영 입대 전 편지", "입대 소식·일정 안내 (이름 빠짐)"]],
      Inches(0.7), Inches(2.5), Inches(12), [Inches(2.4), Inches(4.8), Inches(4.8)], fs=13)
tf = box(s, Inches(0.7), Inches(5.2), Inches(12), Inches(1.8))
setp(tf.paragraphs[0], "2개 탈락 (품질 미달)", 16, ACC, True)
addp(tf, "Gemma-12B → 브랜드명 환각(없는 단어 생성)   /   InternVL3 → 한국어 칸에 중·일 누출(순도 27%)", 14, INK, bullet=True)
addp(tf, "→ 상위 3개: Qwen3.5-9B · Qwen3-VL-8B · Gemma-26B", 15, NAVY, True)

# 7. 4단계 최종
s = slide(); title_bar(s, "STEP 4", "남은 3개 최종 비교 → 선정")
tf = box(s, Inches(0.7), Inches(1.65), Inches(12), Inches(0.6))
setp(tf.paragraphs[0], "평가 지표 직접 설계: 번체율 / 한국어 순도(한자 잔류) 자동 측정", 13, GRAY)
table(s, [["모델", "번체율", "한국어 순도", "속도", "메모리", "판정"],
          ["Gemma-26B", "100%", "100%", "6.6s", "14GB", "품질 1위"],
          ["Qwen3.5-9B", "100%", "93%", "7.0s", "5GB", "실용 1위"],
          ["Qwen3-VL-8B", "100%", "47%", "9.3s", "9GB", "탈락(한자잔류)"]],
      Inches(0.7), Inches(2.4), Inches(12),
      [Inches(2.5), Inches(1.6), Inches(2.4), Inches(1.5), Inches(1.8), Inches(2.2)])
tf = box(s, Inches(0.7), Inches(4.5), Inches(12), Inches(2.5))
setp(tf.paragraphs[0], "최종 선정: Qwen3.5-9B", 20, ACC, True)
addp(tf, "품질 거의 동급(순도 93%)이면서 가장 가볍고(5GB) 빠르며 24GB에 안전 → 배포·운영 현실성", 14, INK, bullet=True)
addp(tf, "Gemma-26B는 순수 품질 최댓값(순도 100%)이나 메모리 14GB 부담 → 차순위", 13, GRAY, bullet=True)

# 8. 인사이트
s = slide(); title_bar(s, "INSIGHTS", "핵심 인사이트")
tf = box(s, Inches(0.7), Inches(2.0), Inches(12), Inches(5))
setp(tf.paragraphs[0], "1.  큰 모델 ≠ 더 좋다 — 27B는 못 돌고(탈락), 큰 모델이 오히려 부정확하기도", 16, INK, True)
addp(tf, "2.  결정 요인은 크기가 아니라 '지시 준수' — 번체/간체, 한자 잔류, 환각 여부", 16, INK, True, space=14)
addp(tf, "3.  같은 모델도 프롬프트·입력방식·처리순서 설계로 품질이 급변", 16, INK, True, space=14)
addp(tf, "4.  평가 지표를 직접 설계(번체율·한국어 순도)해 주관 아닌 정량 비교", 16, INK, True, space=14)
addp(tf, "5.  한계(고유명사 음역)는 RAG로 보완 예정 · fine-tuning은 효과 대비 비용 커 배제", 16, INK, True, space=14)

prs.save("vl_results/발표대본/VLM_모델선정.pptx")
print("저장 완료: vl_results/발표대본/VLM_모델선정.pptx")
print("슬라이드 수:", len(prs.slides._sldIdLst))
